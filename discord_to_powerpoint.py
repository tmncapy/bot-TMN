import discord
from discord.ext import commands
from pptx import Presentation
import os

# Khởi tạo bot với tiền tố !
bot = commands.Bot(command_prefix="!")

# Tạo hoặc mở file PowerPoint macro-enabled
ppt_file = "output.pptm"
if not os.path.exists(ppt_file):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])  # Slide trắng
    prs.save(ppt_file)

@bot.event
async def on_ready():
    print(f'Bot đã đăng nhập thành công với tên {bot.user}')

@bot.event
async def on_message(message):
    if message.author == bot.user:
        return  # Bỏ qua tin nhắn của bot

    # Liên kết kênh Discord với các shape tương ứng
    channel_to_shape = {
        123456789012345678: "a1",  # ID kênh 1 -> shape "a1"
        223456789012345678: "a2",  # ID kênh 2 -> shape "a2"
        323456789012345678: "a3",  # ID kênh 3 -> shape "a3"
        423456789012345678: "a4"   # ID kênh 4 -> shape "a4"
    }

    shape_name = channel_to_shape.get(message.channel.id)
    if not shape_name:
        return  # Bỏ qua nếu kênh không nằm trong danh sách

    # Mở file PowerPoint và cập nhật nội dung vào shape tương ứng
    prs = Presentation(ppt_file)
    slide = prs.slides[0]  # Lấy slide đầu tiên

    # Tìm shape theo tên
    target_shape = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target_shape = shape
            break

    if target_shape and hasattr(target_shape, "text_frame"):
        target_shape.text = message.content  # Ghi nội dung tin nhắn
        prs.save(ppt_file)
        await message.channel.send(f'Nội dung "{message.content}" đã được thêm vào slide 1, shape "{shape_name}"!')
    else:
        await message.channel.send(f'Không tìm thấy shape "{shape_name}" trên slide 1.')

# Chạy bot
TOKEN = "YOUR_DISCORD_BOT_TOKEN"
bot.run(TOKEN)

